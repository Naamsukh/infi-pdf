import os
from io import BytesIO

from pptx.util import Inches, Pt
from pptx import Presentation
from pptx.dml.color import RGBColor

from overlap_utils import adjust_overlapping_objects_in_ppt, find_unique_overlapping_and_non_overlapping_objects

def organize_data_by_page(json_data):
    # Group data by page number
    pages = {}
    for item in json_data:
        page_number = item['metadata'].get('page_number', 1)  # Default to page 1 if not specified
        if page_number not in pages:
            pages[page_number] = []
        pages[page_number].append(item)
    
    # Sort pages and return as a list of lists
    return [pages[key] for key in sorted(pages.keys())]

def create_ppt_demo(json_data, slide_width_inch=13.33, slide_height_inch=7.5):
    prs = Presentation()

    # Set slide dimensions
    prs.slide_width = Inches(slide_width_inch)
    prs.slide_height = Inches(slide_height_inch)

    # Organize data by page
    pages_data = organize_data_by_page(json_data)

    for page_data in pages_data:
            
        print("Processing page data",page_data[0]['metadata'].get('page_number', 1))    
        
        # Get unique overlap and non overlap elements
        overlaps,non_overlaps = find_unique_overlapping_and_non_overlapping_objects(page_data)
        
        # Adjust the coordinates of the overlapping objects
        adjusted_overlaps_objects = adjust_overlapping_objects_in_ppt(overlaps)

        # Merge the adjusting overlapping objects and non overlapping objects
        adjusted_elements = adjusted_overlaps_objects + non_overlaps

        # Get layout dimensions from the first element in page_data
        layout_width_px = adjusted_elements[0]['metadata']['coordinates']['layout_width']
        layout_height_px = adjusted_elements[0]['metadata']['coordinates']['layout_height']

        # Calculate pixel to inch conversion factors
        px_to_inch_x = slide_width_inch / layout_width_px
        px_to_inch_y = slide_height_inch / layout_height_px

        # Add a blank slide layout
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        for element in adjusted_elements:
            text = element['text']
            element_type = element['type']
            coordinates = element['metadata']['coordinates']['points']

            # Calculate bounding box dimensions in inches
            left_px = min(point[0] for point in coordinates)
            top_px = min(point[1] for point in coordinates)
            right_px = max(point[0] for point in coordinates)
            bottom_px = max(point[1] for point in coordinates)

            left_inch = left_px * px_to_inch_x
            top_inch = top_px * px_to_inch_y
            width_inch = (right_px - left_px) * px_to_inch_x
            height_inch = (bottom_px - top_px) * px_to_inch_y

            if element_type == 'Table':
                # Add a colored block (rectangle) for Table type
                shape = slide.shapes.add_shape(
                    1,  # Rectangle shape
                    Inches(left_inch), Inches(top_inch),
                    Inches(width_inch), Inches(height_inch)
                )
                fill = shape.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(192, 192, 192)  # Light gray color for Table blocks
                shape.text_frame.text = "TP"
            elif element_type == 'Image':
                # Add the image to the slide
                image_path = element['metadata'].get('image_path')
                if image_path and os.path.exists(image_path):
                    slide.shapes.add_picture(image_path, Inches(left_inch), Inches(top_inch),
                                             Inches(width_inch), Inches(height_inch))
                else:
                    # If image is not found, add a placeholder shape
                    shape = slide.shapes.add_shape(
                        1,  # Rectangle shape
                        Inches(left_inch), Inches(top_inch),
                        Inches(width_inch), Inches(height_inch)
                    )
                    fill = shape.fill
                    fill.solid()
                    fill.fore_color.rgb = RGBColor(192, 192, 192)  # Light gray color for missing images
                    shape.text_frame.text = "IP"
            else:
                # Add text box for other types (e.g., Title, NarrativeText, ListItem)
                textbox = slide.shapes.add_textbox(Inches(left_inch), Inches(top_inch), Inches(width_inch), Inches(height_inch))
                text_frame = textbox.text_frame
                text_frame.text = text
                text_frame.word_wrap = True

                # Dynamically adjust the font size based on the bounding box
                optimal_font_size = fit_text_in_box(text_frame, width_inch, height_inch)
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(optimal_font_size)
    
    ppt_io = BytesIO()
    prs.save(ppt_io)
    ppt_io.seek(0)  # Reset buffer pointer

    return ppt_io
    # # Save the presentation
    # prs.save('output_presentation_multiple_pages.pptx')

def fit_text_in_box(text_frame, width_inch, height_inch, font_name='Calibri'):
    max_font_size = 36
    min_font_size = 6

    for font_size in range(max_font_size, min_font_size - 1, -1):
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = font_name
                run.font.size = Pt(font_size)

        lines = text_frame.text.count("\n") + 1
        approx_text_height = lines * 1.2 * (font_size / 72)
        approx_text_width = len(text_frame.text) * 0.6 * (font_size / 72)

        if approx_text_height <= height_inch and approx_text_width <= width_inch:
            return font_size  # Return the next smaller font size

    return min_font_size